"""
utils_a.py

This module contains utility functions and constants used across the IIA QA Automation Python Script project. 


=================================================================================================================
A.
Review the Email Library for cases we should log based on Email archive - start w/ attachments sent to HSR in Email Library 

DV sheet is in Working Folder, DV is filled out, copy in Outgoing Folder has no DV sheet 

DV sheet is in Working Folder, DV is filled out, copy in Outfolder Folder has completed DV (forgot to remove it) 

DV sheet is in Working Folder, DV is not filled out, copy in Outgoing Folder has blank DV sheet 

DV sheet is in Working Folder, DV is not filled out, copy in Outgoing Folder has no DV sheet 

No DV sheet in either Working Folder or Outgoing Folder. 

Attachment has no copy in any Outgoing Folder 

Note that many Attachments are screened out 

Recipient must be in hsr.ca.gov domain but NOT one of the internal team members with an hsr.ca.gov email 

Only PDF, PPT and DOCX filetypes - review DOCX and PPT for the DV sheet 

If PDF is attached alone, look for correspoding DOCX or PPT filename in Outgoing and Working folders 


B.
Review all files in Outgoing Folders to see if there is an email in Email Library that has an attachment that matches the filename.  

Create an XL log of Outgoing files with no email . Include names of ppeople who worked on the file. 
=================================================================================================================

Constants:
- TASK_FOLDERS: A list of dictionaries containing information about various task folders, including their names, drive IDs, working folder IDs, and outgoing folder IDs.
- FILE_OUTPUT_NAME: The name of the Excel file where the QA Automation results will be saved.
"""

from cred import *
import requests
from lxml import etree
import zipfile
from pptx import Presentation
import time
import os
import io
import pandas as pd
import extract_msg
import traceback
import re
from auth import *

TASK_FOLDERS = [
    { "folder_name": "Task 01 - Renewable Energy", 
    "task_drive_id": TASK_01_DRIVE_ID, 
    "working_folder_id" : TASK_01_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_01_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 02 - Climate Adaptation", 
    "task_drive_id": TASK_02_DRIVE_ID, 
    "working_folder_id" : TASK_02_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_02_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 03 - Strategy", 
    "task_drive_id": TASK_03_DRIVE_ID, 
    "working_folder_id" : TASK_03_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_03_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 04 - Reporting", 
    "task_drive_id": TASK_04_DRIVE_ID, 
    "working_folder_id" : TASK_04_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_04_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 05 - GHG & Air Quality", 
    "task_drive_id": TASK_05_DRIVE_ID, 
    "working_folder_id" : TASK_05_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_05_OUTGOING_FILES_ID
    },
]

FILE_OUTPUT_NAME = "QA_Automation_Output-test.xlsx"
FILE_OUTPUT_NAME_TASK_B = "QA_Automation_Missing_Email_Attachments_Output.xlsx"
"""
Get the id of the .msg file in Email Library
"""
def get_weburl_item_id(web_url, drives_id,  headers):
    file_name = web_url.split("/")[-1]
    msg_item_url = f'{GRAPH_API_URL}/drives/{drives_id}/root:/{file_name}'
    response = requests.get(msg_item_url, headers=headers)
    item_id = response.json().get('id')
    return item_id

"""
List the children of a SharePoint folder
"""
def list_children(site_id, drive_id, item_id, headers):
    url = f"{GRAPH_API_URL}/sites/{site_id}/drives/{drive_id}/items/{item_id}/children"
    response = requests.get(url, headers=headers)
    return response.json()

"""
Get the content of a file in SharePoint by its ID
"""
def fetch_file_content(site_id, drive_id, file_id, headers):
    url = f"{GRAPH_API_URL}/sites/{site_id}/drives/{drive_id}/items/{file_id}/content"
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.content  # Return raw file content
    else:
        return None

"""
Using search in graphapi to recursively find the partial match of the file name in SharePoint folder
"""
def search_files_in_folder(site_id, drive_id, folder_id, query, headers):
    url = f"{GRAPH_API_URL}/sites/{site_id}/drives/{drive_id}/items/{folder_id}/search(q='{query}')"
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.json()  # Return the search results
    else:
        return None

"""
Check for matching extension for .docx and .pptx from the returned search results
"""
def find_file_in_subfolders(site_id, drive_id, parent_item_id, filename, headers):
    search_results = search_files_in_folder(site_id, drive_id, parent_item_id, filename, headers)
    if search_results:
        # print("THE SEARCH RESULTS")
        # print(search_results)
        valid_extensions = ('.docx', '.pptx')
        for item in search_results['value']:
            # Check if the file name matches and has a valid extension
            if item['name'].lower().startswith(filename.lower()) and item['name'].lower().endswith(valid_extensions):
                return item  # Return the file item if found
    return None  # Return None if the file is not found


# Function to check if a text is likely a name (first letter capitalized and at least two words)
def is_name(text):
    # Check for common names with at least two words, each starting with a capital letter
    return bool(re.match(r'^[A-Z][a-z]+(?: [A-Z][a-z]+)*$', text))

"""
Check for the presence of DV sheet in the DOCX file
"""
def handle_checking_dv_in_docx_file(msg_response, response_dict):
    doc = io.BytesIO(msg_response.content)
    with zipfile.ZipFile(doc) as docx_zip:
        xml_content = docx_zip.read('word/document.xml')
    
    target_fields = [
        "Signature", "Filename", "Description", 
        "Prepared by", "Checked by", "Approved by", "Name"
    ]

    tree = etree.XML(xml_content)
    namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    text_elements = tree.xpath("//w:t", namespaces=namespaces)

    content = [element.text for element in text_elements if element.text]
    # print(content) 

    if "Document Verification" in content:
        print("DV Sheet found in Working Folder")
        response_dict["is_dv_sheet_exists"] = True

        for field in target_fields:
            for idx, text in enumerate(content):
                if field in text:
                    # Look for the next relevant text (skip non-relevant ones like 'Filename' or 'Description')
                    next_text = None
                    for i in range(idx + 1, len(content)):
                        if content[i] not in target_fields and content[i].strip() != "" and content[i] not in ["Filename", "Description"]:
                            next_text = content[i]
                            break

                    # Check if the next text is a likely name
                    if next_text and is_name(next_text):
                        print(f"Detected name after '{field}': {next_text}")
                        response_dict["is_dv_sheet_filled"] = True
                    else:
                        print(f"nothing detected after '{field}'")

    else:
        print("No DV Sheet in Working Folder")



def handle_checking_dv_in_pptx_file(msg_response, response_dict):
    presentation = Presentation(io.BytesIO(msg_response.content))
    all_text = ""

    # Iterate through slides and extract text from all shapes, including tables
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text + "\n"
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        all_text += cell.text + "\n"
    
    target_fields = [
        "Signature", "Filename", "Description", 
        "Prepared by", "Checked by", "Approved by", "Name"
    ]

    
    if "Document Verification" in all_text:
        print("DV Sheet found in Presentation")
        response_dict["is_dv_sheet_exists"] = True

        lines = all_text.split('\n')
        for field in target_fields:
            detected = False
            for i, line in enumerate(lines):
                if field in line:
                    # Look for the next relevant text (skip non-relevant ones like 'Filename' or 'Description')
                    next_text = None
                    for j in range(i + 1, len(lines)):
                        if lines[j].strip() and lines[j].strip() not in target_fields and lines[j].strip() not in ["Filename", "Description"]:
                            next_text = lines[j].strip()
                            break

                    # Check if the next text is a likely name
                    if next_text:
                        if is_name(next_text):
                            print(f"Detected name after '{field}': {next_text}")
                            response_dict["is_dv_sheet_filled"] = True
                            detected = True
                            break

            if not detected:
                print(f"nothing detected after '{field}'")
    else:
        print("No DV Sheet in Presentation")


"""
Returns the following statuses in an dictionary

1. whether file exists
2. whether the DV sheet exists in specified folder
3. whether the DV sheet is filled out or is blank

example result
{
    "is_file_exists": True,
    "is_dv_sheet_exists": True,
    "is_dv_sheet_filled": True
}

"""
def handle_DV_sheet_exists_status(headers, 
                           task_drive_id,
                           task_folder_id,
                           attachment_name_without_extension):
    
    response_dict = {
        "is_file_exists": False,
        "is_dv_sheet_exists": False,
        "is_dv_sheet_filled": False
    }
    file = find_file_in_subfolders(site_id=SITE_ID, drive_id=task_drive_id, parent_item_id=task_folder_id, filename=attachment_name_without_extension, headers=headers)

    if file:
        print(f"File found: {file['name']}")
        response_dict["is_file_exists"] = True
    else:
        print("File not found")
        return response_dict
    
    item_content_url  = f'{GRAPH_API_URL}/drives/{file["parentReference"]["driveId"]}/items/{file["id"]}/content'
    msg_response = requests.get(item_content_url, headers=headers)

    file_extension = os.path.splitext(file['name'])[1].lower()

    if file_extension == '.pptx':
        print("The file is a PPTX.")
        handle_checking_dv_in_pptx_file(msg_response=msg_response, response_dict=response_dict)
    elif file_extension == '.docx':
        print("The file is a DOCX.")
        handle_checking_dv_in_docx_file(msg_response=msg_response, response_dict=response_dict)
    else:
        print("The file is neither a pptx nor a DOCX.")
    
    return response_dict


def get_list_of_internal_members(headers):
    internal_members = []
    drives_url = f'{GRAPH_API_URL}/sites/{SITE_ID}/lists/{INTERNAL_MEMBERS_LIST_ID}/items?$expand=fields'
    response = requests.get(drives_url, headers=headers)
    data = response.json()

    for item in data['value']:
        # TODO: change to @HSR column field once created
        internal_members.append(item['fields']['Title'])
    return internal_members

def process_dv_dataframes(working_folder_statuses, 
                          outgoing_folder_statuses, 
                          attachment_name, 
                          dataframes,
                          msg_info):
    

    df_data = {
        "Email Subject": msg_info['subject'],
        "Sender": msg_info['sender'],
        "Recipients": msg_info['recipients'],
        "Attachment Name": attachment_name,
        "Working Folder Exists": working_folder_statuses['is_file_exists'],
        "Outgoing Folder Exists": outgoing_folder_statuses['is_file_exists']
    }

    # 1. DV sheet is in Working Folder, DV is filled out, copy in Outgoing Folder has no DV sheet 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and working_folder_statuses['is_dv_sheet_filled'] \
        and not outgoing_folder_statuses['is_dv_sheet_exists']:
        dataframes['dv_filled_in_working_and_outgoing_no_dv_sheet_df'] = dataframes['dv_filled_in_working_and_outgoing_no_dv_sheet_df'].append(
            df_data, ignore_index=True)
        
    # 2. DV sheet is in Working Folder, DV is filled out, copy in Outfolder Folder has completed DV (forgot to remove it) 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and working_folder_statuses['is_dv_sheet_filled'] \
        and outgoing_folder_statuses['is_dv_sheet_exists'] and outgoing_folder_statuses['is_dv_sheet_filled']:
        dataframes['dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df'] = dataframes['dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df'].append(
            df_data, ignore_index=True)
    
    # 3. DV sheet is in Working Folder, DV is not filled out, copy in Outgoing Folder has blank DV sheet 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and not working_folder_statuses['is_dv_sheet_filled'] \
        and outgoing_folder_statuses['is_dv_sheet_exists'] and not outgoing_folder_statuses['is_dv_sheet_filled']:
        dataframes['dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df'] = dataframes['dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df'].append(
            df_data, ignore_index=True)
        
    # 4. DV sheet is in Working Folder, DV is not filled out, copy in Outgoing Folder has no DV sheet 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and not working_folder_statuses['is_dv_sheet_filled'] \
        and outgoing_folder_statuses['is_file_exists'] and not outgoing_folder_statuses['is_dv_sheet_exists']:
        dataframes['dv_not_filled_in_working_outgoing_no_dv_sheet_df'] = dataframes['dv_not_filled_in_working_outgoing_no_dv_sheet_df'].append(
            df_data, ignore_index=True)
    
    # 5. No DV sheet in either Working Folder or Outgoing Folder.  
    if not working_folder_statuses['is_file_exists'] or not outgoing_folder_statuses['is_file_exists']:
        dataframes['no_dv_in_working_or_outgoing_df'] = dataframes['no_dv_in_working_or_outgoing_df'].append(
            df_data, ignore_index=True)

    # 6. Attachment has no copy in any Outgoing Folder 
    if not outgoing_folder_statuses['is_file_exists']:
        dataframes['attachment_not_found_in_outgoing_folder_df'] = dataframes['attachment_not_found_in_outgoing_folder_df'].append(
            df_data, ignore_index=True)
        
"""
Run the QA Automation processing in the background for task A
"""
def run_qa_automation_A_in_background(drives_url, headers, refresh_token=None):
    print("Running QA Automation A in the background...")
    start_time = time.time()
    run_qa_automation_processing(drives_url=drives_url, headers=headers, refresh_token=refresh_token)
    end_time = time.time()
    execution_time = end_time - start_time
    print(f"Execution time: {execution_time} seconds for QA Automation A")

"""
Run the QA Automation processing in the background for task B
"""
def run_qa_automation_B_in_background(drives_url, headers, refresh_token):
    print("Running QA Automation B in the background...")
    start_time = time.time()
    print("running task B")
    run_qa_automation_processing_B(drives_url=drives_url, headers=headers, refresh_token=refresh_token)
    print("finished task B")

    end_time = time.time()
    execution_time = end_time - start_time
    print(f"Execution time: {execution_time} seconds for QA Automation B")

def run_qa_automation_processing_B(drives_url, headers, refresh_token=None):

    dataframes = {
        'files_not_found_in_email_library': pd.DataFrame()
    }
    response = requests.get(drives_url, headers=headers)
    data = response.json().get('value', [])
    valid_extensions = ('.docx', '.pptx', '.pdf')

    print("data length", len(data))
    filtered_data = [
        item for item in data
        if '@HSR' in item.get('fields', {}).get('Arup_To', '')
    ]

    result_filted_data = []
    internal_members = get_list_of_internal_members(headers)

    for item in filtered_data:
        email = item['fields'].get('Arup_To')
        if not email:
            continue

        name = switch_name_format(email)
        if name:
            if name not in internal_members:
                result_filted_data.append(item)

    # List to store error logs
    error_logs = []
    attachments_from_email_library = []


    print("LENGHT OF RESULT FILTERED DATA for non internal members", len(result_filted_data))


    for i, item in enumerate(result_filted_data):
        print("====Item processing===", i)
        valid_extensions = ('.docx', '.pptx', '.pdf')
        
        # Get refresh token after every 30 items
        if i >= 30 and i % 30 == 0:
            refresh_token_updated, expires_in = refresh_user_token(refresh_token)
            print("Renerated new token")
            headers = {'Authorization': 'Bearer ' + refresh_token_updated}

        try:
            web_url = item['webUrl']
            item_id = get_weburl_item_id(web_url=web_url, drives_id=EMAIL_LIBRARY_DRIVE_ID,headers=headers)
            item_content_url  = f'{GRAPH_API_URL}/drives/{EMAIL_LIBRARY_DRIVE_ID}/items/{item_id}/content'
            msg_response = requests.get(item_content_url, headers=headers)
            msg = extract_msg.Message(io.BytesIO(msg_response.content))
            attachments = msg.attachments

            for attachment in attachments:

                if not attachment.longFilename:
                        print("No filename skipping")
                        continue
                
                if not attachment.longFilename.lower().endswith(valid_extensions):
                        print("not valid ext skipping: " + attachment.longFilename)
                        continue
                attachment_name = attachment.longFilename
                attachment_name_without_extension = os.path.splitext(attachment_name)[0]
                attachments_from_email_library.append(attachment_name_without_extension)


        except Exception as e:
            print("An error occured processing: ", e)
            traceback.print_exc()
            log_error(item=item, error_message=e, error_logs=error_logs)
    
    print("ATTACHMENTS FROM EMAIL LIBRARY")
    print(attachments_from_email_library[:10])
    print("LENGTH OF ATTACHMENTS FROM EMAIL LIBRARY", len(attachments_from_email_library))
    print("\n\n===FINDING ATTACHMENTS IN OUTGOING FOLDERS TO MATCH WITH EMAIL LIBRARY===")
    for i, task_folder in enumerate(TASK_FOLDERS):

        try:
            task_drive_id = task_folder['task_drive_id']
            task_outgoing_folder_id = task_folder['outgoing_folder_id']

            search_results_docx = search_files_in_folder(site_id=SITE_ID, 
                                    drive_id=task_drive_id, 
                                    folder_id=task_outgoing_folder_id, 
                                    query='.docx', headers=headers
                                    )
            
            if search_results_docx:
                for item in search_results_docx['value']:

                    file_name_without_extension = os.path.splitext(item['name'])[0]
                    created_by_name = item['createdBy']['user']['displayName']
                    created_by_email = item['createdBy']['user']['email']
                    last_modified_by_name = item['lastModifiedBy']['user']['displayName']
                    last_modified_by_email = item['lastModifiedBy']['user']['email']

                    if file_name_without_extension not in attachments_from_email_library:
                        print("Attachment not found in Email Library: ", item['name'])
                        df_data = {
                            "Attachment Name": item['name'],
                            "Created By Name": created_by_name,
                            "Created By Email": created_by_email,
                            "Last Modified By Name": last_modified_by_name,
                            "Last Modified By Email": last_modified_by_email
                        }
                        dataframes['files_not_found_in_email_library'] = dataframes['files_not_found_in_email_library'].append(
                            df_data, ignore_index=True)

                    
            
            search_results_pptx = search_files_in_folder(site_id=SITE_ID, 
                                    drive_id=task_drive_id, 
                                    folder_id=task_outgoing_folder_id, 
                                    query='.pptx', headers=headers
                                    )

            if search_results_pptx:
                for item in search_results_pptx['value']:

                    file_name_without_extension = os.path.splitext(item['name'])[0]
                    created_by_name = item['createdBy']['user']['displayName']
                    created_by_email = item['createdBy']['user']['email']
                    last_modified_by_name = item['lastModifiedBy']['user']['displayName']
                    last_modified_by_email = item['lastModifiedBy']['user']['email']

                    if file_name_without_extension not in attachments_from_email_library:
                        print("Attachment not found in Email Library: ", item['name'])
                        df_data = {
                            "Attachment Name": item['name'],
                            "Created By Name": created_by_name,
                            "Created By Email": created_by_email,
                            "Last Modified By Name": last_modified_by_name,
                            "Last Modified By Email": last_modified_by_email
                        }

                        dataframes['files_not_found_in_email_library'] = dataframes['files_not_found_in_email_library'].append(
                            df_data, ignore_index=True)

        except Exception as e:
            print("An error occured processing: ", e)
            traceback.print_exc()
            log_error(item=item, error_message=e, error_logs=error_logs)
    
                    
    dataframes['files_not_found_in_email_library'].to_excel(FILE_OUTPUT_NAME_TASK_B, sheet_name='Attachment Not Found', index=False)
    print("Excel file created successfully for task B") 



# Function to log errors
def log_error(item, error_message, error_logs):
    msg_info = {
        "item": item,
        "error_message": error_message
    }
    error_logs.append(msg_info)


# Function to loop through all sheets and add items in the 'attachments' column to an array
def get_attachments_from_excel(file_path):
    attachments = []
    try:
        # Load the Excel workbook
        xls = pd.ExcelFile(file_path)
        # Loop through each sheet in the workbook
        for sheet_name in xls.sheet_names:
            # Read the sheet into a DataFrame
            df = pd.read_excel(xls, sheet_name=sheet_name)
            
            # Check if the 'attachments' column exists in the DataFrame
            if 'Attachment Name' in df.columns:
                # Add the items in the 'attachments' column to the list
                attachments.extend(df['Attachment Name'].dropna().tolist())
        
        return attachments
    except Exception as e:
        return attachments


"""
Run the QA Automation processing
"""
def run_qa_automation_processing(drives_url, headers, refresh_token=None):
    dataframes = {
        'dv_filled_in_working_and_outgoing_no_dv_sheet_df': pd.DataFrame(),
        'dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df': pd.DataFrame(),
        'dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df': pd.DataFrame(),
        'dv_not_filled_in_working_outgoing_no_dv_sheet_df': pd.DataFrame(),
        'no_dv_in_working_or_outgoing_df': pd.DataFrame(),
        'attachment_not_found_in_outgoing_folder_df': pd.DataFrame()
    }

    # Define the sheet names for each DataFrame
    sheet_names = {
        'dv_filled_in_working_and_outgoing_no_dv_sheet_df': '1',
        'dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df': '2',
        'dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df': '3',
        'dv_not_filled_in_working_outgoing_no_dv_sheet_df': '4',
        'no_dv_in_working_or_outgoing_df': '5',
        'attachment_not_found_in_outgoing_folder_df': '6'
    }

    response = requests.get(drives_url, headers=headers)

    data = response.json().get('value', [])

    print("data length", len(data))
    filtered_data = [
        item for item in data
        if '@HSR' in item.get('fields', {}).get('Arup_To', '')
    ]

    result_filted_data = []
    internal_members = get_list_of_internal_members(headers)

    for item in filtered_data:
        email = item['fields'].get('Arup_To')
        if not email:
            continue

        name = switch_name_format(email)
        if name:
            if name not in internal_members:
                result_filted_data.append(item)

    # List to store error logs
    error_logs = []

    print("LENGHT OF RESULT FILTERED DATA for non internal members", len(result_filted_data))

    # Retrieve previously process attachments from the output file if applicable to speed up processing by skipping the ones in the list.
    cached_attachements_from_previous_output = get_attachments_from_excel(FILE_OUTPUT_NAME)

    for i, item in enumerate(result_filted_data):
        print("====Item processing===", i)
        valid_extensions = ('.docx', '.pptx', '.pdf')
        
        # Get refresh token after every 30 items
        if i >= 30 and i % 30 == 0:
            refresh_token_updated, expires_in = refresh_user_token(refresh_token)
            print("Renerated new token")
            headers = {'Authorization': 'Bearer ' + refresh_token_updated}

        # DEBUG REMOVE LATER. PURPOSE to only test specific items
        # if i != 11 and i != 16 and i != 17:
        #     continue

        try:
            web_url = item['webUrl']
            item_id = get_weburl_item_id(web_url=web_url, drives_id=EMAIL_LIBRARY_DRIVE_ID,headers=headers)
            item_content_url  = f'{GRAPH_API_URL}/drives/{EMAIL_LIBRARY_DRIVE_ID}/items/{item_id}/content'
            msg_response = requests.get(item_content_url, headers=headers)
            msg = extract_msg.Message(io.BytesIO(msg_response.content))
            attachments = msg.attachments

            msg_info = {
                "subject": msg.subject,
                "sender": msg.sender,
                "recipients": msg.to,
                "attachment_count": len(msg.attachments)
            }

            for attachment in attachments:

                if "Acceptable_Use_Acknowledgement" in attachment.longFilename or \
                "Acceptable Use Acknowledgement" in attachment.longFilename:
                    print("Acceptable_Use_Acknowledgement skipping")
                    continue

                if attachment.longFilename in cached_attachements_from_previous_output:
                    print("Attachment already cached skipping")
                    continue

                if not attachment.longFilename:
                    print("No filename skipping")
                    continue

                if not attachment.longFilename.lower().endswith(valid_extensions):
                    print("not valid ext skipping: " + attachment.longFilename)
                    continue

                print(f"Attachment: {attachment.longFilename}")
                attachment_name = attachment.longFilename
                attachment_name_without_extension = os.path.splitext(attachment_name)[0]
                attachment_name_without_extension = attachment_name_without_extension.strip()

                for i, task_folder in enumerate(TASK_FOLDERS):
                    task_drive_id = task_folder['task_drive_id']
                    task_working_folder_id = task_folder['working_folder_id']
                    
                    print("Scanning working folder")
                    working_folder_statuses = handle_DV_sheet_exists_status(headers=headers,
                                           task_drive_id=task_drive_id,
                                           task_folder_id=task_working_folder_id,
                                           attachment_name_without_extension=attachment_name_without_extension)    

                    if working_folder_statuses['is_file_exists']:
                        print("Working folder file exists")
                        break

                for i, task_folder in enumerate(TASK_FOLDERS):
                    task_drive_id = task_folder['task_drive_id']
                    task_outgoing_folder_id = task_folder['outgoing_folder_id']

                    print("Scanning outgoing folder")
                    outgoing_folder_statuses = handle_DV_sheet_exists_status(
                        headers=headers,
                        task_drive_id=task_drive_id,
                        task_folder_id=task_outgoing_folder_id,
                        attachment_name_without_extension=attachment_name_without_extension) 

                    if outgoing_folder_statuses['is_file_exists']:
                        print("Outgoing folder file exists")
                        break
                
                print("==========RESULTS===========")
                print("Working folder statuses")
                print(working_folder_statuses)
                print("Outgoing folder statuses")
                print(outgoing_folder_statuses)

                process_dv_dataframes(working_folder_statuses=working_folder_statuses,
                                    outgoing_folder_statuses=outgoing_folder_statuses, 
                                    attachment_name=attachment_name, 
                                    dataframes=dataframes,
                                    msg_info=msg_info)

        except Exception as e:
            print("An error occured processing: ", e)
            traceback.print_exc()
            log_error(item=item, error_message=e, error_logs=error_logs)
            
    

    # Loop through each DataFrame and add the new data
    for df_name, sheet_name in sheet_names.items():
        try:
            # Try to read the existing sheet into a DataFrame
            existing_df = pd.read_excel(FILE_OUTPUT_NAME, sheet_name=sheet_name)
        except Exception as e:
            # If the sheet doesn't exist, create an empty DataFrame
            existing_df = pd.DataFrame()
        
        # Combine the existing DataFrame with the new data
        dataframes[df_name] = pd.concat([existing_df, dataframes[df_name]], ignore_index=True)

    # with pd.ExcelWriter(FILE_OUTPUT_NAME, engine="xlsxwriter") as writer:
    #     for df_name, df in dataframes.items():
    #         sheet_name = sheet_names[df_name]
    #         df.to_excel(writer, sheet_name=sheet_name, index=False)

    # Check if the file exists
    if not os.path.exists(FILE_OUTPUT_NAME):
        # Create a new Excel file
        with pd.ExcelWriter(FILE_OUTPUT_NAME, engine='openpyxl') as writer:
            for df_name, sheet_name in sheet_names.items():
                dataframes[df_name].to_excel(writer, sheet_name=sheet_name, index=False)
    else:
        # Append to the existing Excel file
        with pd.ExcelWriter(FILE_OUTPUT_NAME, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
            for df_name, sheet_name in sheet_names.items():
                dataframes[df_name].to_excel(writer, sheet_name=sheet_name, index=False)
            
    print(F"Process completed. Results saved to {FILE_OUTPUT_NAME}")
    df_error_logs = pd.DataFrame(error_logs)



    with pd.ExcelWriter(f'Errors_for_{FILE_OUTPUT_NAME}', engine="xlsxwriter") as writer:
        df_error_logs.to_excel(writer, sheet_name='Error Logs', index=False)

    print(f"Error logs saved to Errors_for_{FILE_OUTPUT_NAME}")
    return "Completed processing"


"""
Function to switch names from "Lastname, First@HSR" to "First Last"
Won't be needed once the SP Project Staff list column has HSR emails to compare against the ones in Email Library
"""
def switch_name_format(email):
    match = re.match(r"(\w+), (\w+)(?:\(.+\))?@HSR", email.strip())
    if match:
        return f"{match.group(2)} {match.group(1)}"
    return None