from cred import *
import requests
from lxml import etree
import zipfile
from pptx import Presentation
import time
import os
import io
import pandas as pd
import extract_msg
import traceback
import re

"""
utils.py

This module contains utility functions and constants used across the IIA QA Automation Python Script project. 

Constants:
- TASK_FOLDERS: A list of dictionaries containing information about various task folders, including their names, drive IDs, working folder IDs, and outgoing folder IDs.
- FILE_OUTPUT_NAME: The name of the Excel file where the QA Automation results will be saved.
"""

TASK_FOLDERS = [
    { "folder_name": "Task 01 - Renewable Energy", 
    "task_drive_id": TASK_01_DRIVE_ID, 
    "working_folder_id" : TASK_01_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_01_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 02 - Climate Adaptation", 
    "task_drive_id": TASK_02_DRIVE_ID, 
    "working_folder_id" : TASK_02_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_02_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 03 - Strategy", 
    "task_drive_id": TASK_03_DRIVE_ID, 
    "working_folder_id" : TASK_03_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_03_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 04 - Reporting", 
    "task_drive_id": TASK_04_DRIVE_ID, 
    "working_folder_id" : TASK_04_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_04_OUTGOING_FILES_ID
    },
    { "folder_name": "Task 05 - GHG & Air Quality", 
    "task_drive_id": TASK_05_DRIVE_ID, 
    "working_folder_id" : TASK_05_WORKING_FILES_ID,
    "outgoing_folder_id": TASK_05_OUTGOING_FILES_ID
    },
]

FILE_OUTPUT_NAME = "QA_Automation_Output_test.xlsx"
"""
Get the id of the .msg file in Email Library
"""
def get_weburl_item_id(web_url, drives_id,  headers):
    file_name = web_url.split("/")[-1]
    msg_item_url = f'{GRAPH_API_URL}/drives/{drives_id}/root:/{file_name}'
    response = requests.get(msg_item_url, headers=headers)
    item_id = response.json().get('id')
    print("THE ITEM ID")
    print(item_id)
    return item_id

"""
List the children of a SharePoint folder
"""
def list_children(site_id, drive_id, item_id, headers):
    url = f"{GRAPH_API_URL}/sites/{site_id}/drives/{drive_id}/items/{item_id}/children"
    response = requests.get(url, headers=headers)
    return response.json()

"""
Get the content of a file in SharePoint by its ID
"""
def fetch_file_content(site_id, drive_id, file_id, headers):
    url = f"{GRAPH_API_URL}/sites/{site_id}/drives/{drive_id}/items/{file_id}/content"
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.content  # Return raw file content
    else:
        return None

"""
Using search in graphapi to recursively find the partial match of the file name in SharePoint folder
"""
def search_files_in_folder(site_id, drive_id, folder_id, query, headers):
    url = f"{GRAPH_API_URL}/sites/{site_id}/drives/{drive_id}/items/{folder_id}/search(q='{query}')"
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.json()  # Return the search results
    else:
        return None

"""
Check for matching extension for .docx and .pptx from the returned search results
"""
def find_file_in_subfolders(site_id, drive_id, parent_item_id, filename, headers):
    search_results = search_files_in_folder(site_id, drive_id, parent_item_id, filename, headers)
    if search_results:
        # print("THE SEARCH RESULTS")
        # print(search_results)
        valid_extensions = ('.docx', '.pptx')
        for item in search_results['value']:
            # Check if the file name matches and has a valid extension
            if item['name'].lower().startswith(filename.lower()) and item['name'].lower().endswith(valid_extensions):
                return item  # Return the file item if found
    return None  # Return None if the file is not found


# Function to check if a text is likely a name (first letter capitalized and at least two words)
def is_name(text):
    # Check for common names with at least two words, each starting with a capital letter
    return bool(re.match(r'^[A-Z][a-z]+(?: [A-Z][a-z]+)*$', text))

"""
Check for the presence of DV sheet in the DOCX file
"""
def handle_checking_dv_in_docx_file(msg_response, response_dict):
    doc = io.BytesIO(msg_response.content)
    with zipfile.ZipFile(doc) as docx_zip:
        xml_content = docx_zip.read('word/document.xml')
    
    target_fields = [
        "Signature", "Filename", "Description", 
        "Prepared by", "Checked by", "Approved by", "Name"
    ]

    tree = etree.XML(xml_content)
    namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    text_elements = tree.xpath("//w:t", namespaces=namespaces)

    content = [element.text for element in text_elements if element.text]
    # print(content) 

    if "Document Verification" in content:
        print("DV Sheet found in Working Folder")
        response_dict["is_dv_sheet_exists"] = True

        for field in target_fields:
            for idx, text in enumerate(content):
                if field in text:
                    # Look for the next relevant text (skip non-relevant ones like 'Filename' or 'Description')
                    next_text = None
                    for i in range(idx + 1, len(content)):
                        if content[i] not in target_fields and content[i].strip() != "" and content[i] not in ["Filename", "Description"]:
                            next_text = content[i]
                            break

                    # Check if the next text is a likely name
                    if next_text and is_name(next_text):
                        print(f"Detected name after '{field}': {next_text}")
                        response_dict["is_dv_sheet_filled"] = True
                    else:
                        print(f"nothing detected after '{field}'")

    else:
        print("No DV Sheet in Working Folder")



def handle_checking_dv_in_pptx_file(msg_response, response_dict):
    presentation = Presentation(io.BytesIO(msg_response.content))
    all_text = ""

    # Iterate through slides and extract text from all shapes, including tables
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text + "\n"
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        all_text += cell.text + "\n"
    
    target_fields = [
        "Signature", "Filename", "Description", 
        "Prepared by", "Checked by", "Approved by", "Name"
    ]

    
    if "Document Verification" in all_text:
        print("DV Sheet found in Presentation")
        response_dict["is_dv_sheet_exists"] = True

        lines = all_text.split('\n')
        for field in target_fields:
            detected = False
            for i, line in enumerate(lines):
                if field in line:
                    # Look for the next relevant text (skip non-relevant ones like 'Filename' or 'Description')
                    next_text = None
                    for j in range(i + 1, len(lines)):
                        if lines[j].strip() and lines[j].strip() not in target_fields and lines[j].strip() not in ["Filename", "Description"]:
                            next_text = lines[j].strip()
                            break

                    # Check if the next text is a likely name
                    if next_text:
                        if is_name(next_text):
                            print(f"Detected name after '{field}': {next_text}")
                            response_dict["is_dv_sheet_filled"] = True
                            detected = True
                            break

            if not detected:
                print(f"nothing detected after '{field}'")
    else:
        print("No DV Sheet in Presentation")


"""
Returns the following statuses in an dictionary

1. whether file exists
2. whether the DV sheet exists in specified folder
3. whether the DV sheet is filled out or is blank

example result
{
    "is_file_exists": True,
    "is_dv_sheet_exists": True,
    "is_dv_sheet_filled": True
}

"""
def handle_DV_sheet_exists_status(headers, 
                           task_drive_id,
                           task_folder_id,
                           attachment_name_without_extension):
    
    response_dict = {
        "is_file_exists": False,
        "is_dv_sheet_exists": False,
        "is_dv_sheet_filled": False
    }
    file = find_file_in_subfolders(site_id=SITE_ID, drive_id=task_drive_id, parent_item_id=task_folder_id, filename=attachment_name_without_extension, headers=headers)

    if file:
        print(f"File found: {file['name']}")
        response_dict["is_file_exists"] = True
    else:
        print("File not found")
        return response_dict
    
    item_content_url  = f'{GRAPH_API_URL}/drives/{file["parentReference"]["driveId"]}/items/{file["id"]}/content'
    msg_response = requests.get(item_content_url, headers=headers)

    file_extension = os.path.splitext(file['name'])[1].lower()

    if file_extension == '.pptx':
        print("The file is a PPTX.")
        handle_checking_dv_in_pptx_file(msg_response=msg_response, response_dict=response_dict)
    elif file_extension == '.docx':
        print("The file is a DOCX.")
        handle_checking_dv_in_docx_file(msg_response=msg_response, response_dict=response_dict)
    else:
        print("The file is neither a pptx nor a DOCX.")
    
    return response_dict


def get_list_of_internal_members(headers):
    internal_members = []
    drives_url = f'{GRAPH_API_URL}/sites/{SITE_ID}/lists/{INTERNAL_MEMBERS_LIST_ID}/items?$expand=fields'
    response = requests.get(drives_url, headers=headers)
    data = response.json()

    for item in data['value']:
        # TODO: change to @HSR column field once created
        internal_members.append(item['fields']['Title'])
    return internal_members

def process_dv_dataframes(working_folder_statuses, 
                          outgoing_folder_statuses, 
                          attachment_name, 
                          dataframes,
                          msg_info):
    

    df_data = {
        "Email Subject": msg_info['subject'],
        "Sender": msg_info['sender'],
        "Recipients": msg_info['recipients'],
        "Attachment Name": attachment_name,
        "Working Folder Exists": working_folder_statuses['is_file_exists'],
        "Outgoing Folder Exists": outgoing_folder_statuses['is_file_exists']
    }

    # 1. DV sheet is in Working Folder, DV is filled out, copy in Outgoing Folder has no DV sheet 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and working_folder_statuses['is_dv_sheet_filled'] \
        and not outgoing_folder_statuses['is_dv_sheet_exists']:
        dataframes['dv_filled_in_working_and_outgoing_no_dv_sheet_df'] = dataframes['dv_filled_in_working_and_outgoing_no_dv_sheet_df'].append(
            df_data, ignore_index=True)
        
    # 2. DV sheet is in Working Folder, DV is filled out, copy in Outfolder Folder has completed DV (forgot to remove it) 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and working_folder_statuses['is_dv_sheet_filled'] \
        and outgoing_folder_statuses['is_dv_sheet_exists'] and outgoing_folder_statuses['is_dv_sheet_filled']:
        dataframes['dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df'] = dataframes['dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df'].append(
            df_data, ignore_index=True)
    
    # 3. DV sheet is in Working Folder, DV is not filled out, copy in Outgoing Folder has blank DV sheet 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and not working_folder_statuses['is_dv_sheet_filled'] \
        and outgoing_folder_statuses['is_dv_sheet_exists'] and not outgoing_folder_statuses['is_dv_sheet_filled']:
        dataframes['dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df'] = dataframes['dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df'].append(
            df_data, ignore_index=True)
        
    # 4. DV sheet is in Working Folder, DV is not filled out, copy in Outgoing Folder has no DV sheet 
    if working_folder_statuses['is_file_exists'] and working_folder_statuses['is_dv_sheet_exists'] and not working_folder_statuses['is_dv_sheet_filled'] \
        and outgoing_folder_statuses['is_file_exists'] and not outgoing_folder_statuses['is_dv_sheet_exists']:
        dataframes['dv_not_filled_in_working_outgoing_no_dv_sheet_df'] = dataframes['dv_not_filled_in_working_outgoing_no_dv_sheet_df'].append(
            df_data, ignore_index=True)
    
    # 5. No DV sheet in either Working Folder or Outgoing Folder.  
    if not working_folder_statuses['is_file_exists'] or not outgoing_folder_statuses['is_file_exists']:
        dataframes['no_dv_in_working_or_outgoing_df'] = dataframes['no_dv_in_working_or_outgoing_df'].append(
            df_data, ignore_index=True)

    # 6. Attachment has no copy in any Outgoing Folder 
    if not outgoing_folder_statuses['is_file_exists']:
        dataframes['attachment_not_found_in_outgoing_folder_df'] = dataframes['attachment_not_found_in_outgoing_folder_df'].append(
            df_data, ignore_index=True)
        
"""
Run the QA Automation processing in the background
"""

def run_qa_automation_in_background(drives_url, headers):
    start_time = time.time()
    run_qa_automation_processing(drives_url=drives_url, headers=headers)
    end_time = time.time()
    execution_time = end_time - start_time
    print(f"Execution time: {execution_time} seconds")

"""
Run the QA Automation processing
"""
def run_qa_automation_processing(drives_url, headers):

    dataframes = {
        'dv_filled_in_working_and_outgoing_no_dv_sheet_df': pd.DataFrame(),
        'dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df': pd.DataFrame(),
        'dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df': pd.DataFrame(),
        'dv_not_filled_in_working_outgoing_no_dv_sheet_df': pd.DataFrame(),
        'no_dv_in_working_or_outgoing_df': pd.DataFrame(),
        'attachment_not_found_in_outgoing_folder_df': pd.DataFrame()
    }

    # Define the sheet names for each DataFrame
    sheet_names = {
        'dv_filled_in_working_and_outgoing_no_dv_sheet_df': '1',
        'dv_filled_in_working_and_outgoing_has_filled_dv_sheet_df': '2',
        'dv_not_filled_in_working_outgoing_no_filled_dv_sheet_df': '3',
        'dv_not_filled_in_working_outgoing_no_dv_sheet_df': '4',
        'no_dv_in_working_or_outgoing_df': '5',
        'attachment_not_found_in_outgoing_folder_df': '6'
    }

    response = requests.get(drives_url, headers=headers)

    data = response.json().get('value', [])

    print("data length", len(data))
    filtered_data = [
        item for item in data
        if '@HSR' in item.get('fields', {}).get('Arup_To', '')
    ]

    result_filted_data = []
    internal_members = get_list_of_internal_members(headers)

    for item in filtered_data:
        email = item['fields'].get('Arup_To')
        if not email:
            continue

        name = switch_name_format(email)
        if name:
            if name not in internal_members:
                result_filted_data.append(item)

    print("LENGHT OF RESULT FILTERED DATA for non internal members", len(result_filted_data))
    for i, item in enumerate(result_filted_data[230:232]):
        print("====Item processing===", i)
        valid_extensions = ('.docx', '.pptx', '.pdf')


        # DEBUG REMOVE LATER. PURPOSE to only test specific items
        # if i != 11 and i != 16 and i != 17:
        #     continue

        try:
            web_url = item['webUrl']
            item_id = get_weburl_item_id(web_url=web_url, drives_id=EMAIL_LIBRARY_DRIVE_ID,headers=headers)
            item_content_url  = f'{GRAPH_API_URL}/drives/{EMAIL_LIBRARY_DRIVE_ID}/items/{item_id}/content'
            msg_response = requests.get(item_content_url, headers=headers)
            msg = extract_msg.Message(io.BytesIO(msg_response.content))
            attachments = msg.attachments

            msg_info = {
                "subject": msg.subject,
                "sender": msg.sender,
                "recipients": msg.to,
                "attachment_count": len(msg.attachments)
            }

            print("Attachments")

            for attachment in attachments:

                if "Acceptable_Use_Acknowledgement" in attachment.longFilename or "Acceptable Use Acknowledgement" in attachment.longFilename:
                    print("Acceptable_Use_Acknowledgement skipping")
                    continue

                print("THE ATTACHEMNT: ", attachment.longFilename)

                if not attachment.longFilename:
                    print("No filename skipping")
                    continue

                if not attachment.longFilename.lower().endswith(valid_extensions):
                    print("not valid ext skipping: " + attachment.longFilename)
                    continue

                print(f"Attachment: {attachment.longFilename}")
                attachment_name = attachment.longFilename
                attachment_name_without_extension = os.path.splitext(attachment_name)[0]
                attachment_name_without_extension = attachment_name_without_extension.strip()

                print("FINDING:", attachment_name_without_extension)

                for i, task_folder in enumerate(TASK_FOLDERS):
                    task_drive_id = task_folder['task_drive_id']
                    task_working_folder_id = task_folder['working_folder_id']
                    
                    print("Scanning working folder")
                    working_folder_statuses = handle_DV_sheet_exists_status(headers=headers,
                                           task_drive_id=task_drive_id,
                                           task_folder_id=task_working_folder_id,
                                           attachment_name_without_extension=attachment_name_without_extension)    

                    if working_folder_statuses['is_file_exists']:
                        print("Working folder file exists")
                        break

                for i, task_folder in enumerate(TASK_FOLDERS):
                    task_drive_id = task_folder['task_drive_id']
                    task_outgoing_folder_id = task_folder['outgoing_folder_id']

                    print("Scanning outgoing folder")
                    outgoing_folder_statuses = handle_DV_sheet_exists_status(
                        headers=headers,
                        task_drive_id=task_drive_id,
                        task_folder_id=task_outgoing_folder_id,
                        attachment_name_without_extension=attachment_name_without_extension) 

                    if outgoing_folder_statuses['is_file_exists']:
                        print("Outgoing folder file exists")
                        break
                
                print("==========RESULTS===========")
                print("Working folder statuses")
                print(working_folder_statuses)
                print("Outgoing folder statuses")
                print(outgoing_folder_statuses)

                process_dv_dataframes(working_folder_statuses=working_folder_statuses,
                                    outgoing_folder_statuses=outgoing_folder_statuses, 
                                    attachment_name=attachment_name, 
                                    dataframes=dataframes,
                                    msg_info=msg_info)

        except Exception as e:
            print("An error occured processing: ", e)
            traceback.print_exc()
            
            

    with pd.ExcelWriter(FILE_OUTPUT_NAME, engine="xlsxwriter") as writer:
        for df_name, df in dataframes.items():
            sheet_name = sheet_names[df_name]
            df.to_excel(writer, sheet_name=sheet_name, index=False)

    print(F"Process completed. Results saved to {FILE_OUTPUT_NAME}")
    return filtered_data


"""
Function to switch names from "Lastname, First@HSR" to "First Last"
Won't be needed once the SP Project Staff list column has HSR emails to compare against the ones in Email Library
"""
def switch_name_format(email):
    match = re.match(r"(\w+), (\w+)(?:\(.+\))?@HSR", email.strip())
    if match:
        return f"{match.group(2)} {match.group(1)}"
    return None